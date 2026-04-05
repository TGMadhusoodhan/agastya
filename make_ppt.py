from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Colour palette ─────────────────────────────────────────────────────────────
DEEP_NAVY    = RGBColor(0x0A, 0x0F, 0x2E)   # slide background
TEAL         = RGBColor(0x00, 0xC9, 0xB1)   # accent / headings
GOLD         = RGBColor(0xFF, 0xC0, 0x36)   # highlight
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCC, 0xD6, 0xE8)
CARD_BG      = RGBColor(0x12, 0x1A, 0x42)   # card fill

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color=DEEP_NAVY):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill_color=CARD_BG, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def accent_bar(slide, l, t, w=0.07, h=0.5, color=TEAL):
    bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def bullet_box(slide, items, l, t, w, h, size=15, color=WHITE, dot_color=TEAL, gap=0.38):
    for i, item in enumerate(items):
        y = t + i * gap
        # dot
        d = slide.shapes.add_shape(1, Inches(l), Inches(y + 0.10), Inches(0.09), Inches(0.09))
        d.fill.solid(); d.fill.fore_color.rgb = dot_color; d.line.fill.background()
        # text
        tb = slide.shapes.add_textbox(Inches(l + 0.18), Inches(y), Inches(w - 0.18), Inches(gap))
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run(); run.text = item
        run.font.size = Pt(size); run.font.color.rgb = color

def divider(slide, t, color=TEAL, l=0.5, w=12.33):
    bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.03))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()

def tag_pill(slide, label, l, t, bg_color=TEAL, text_color=DEEP_NAVY):
    b = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(len(label)*0.11+0.25), Inches(0.33))
    b.fill.solid(); b.fill.fore_color.rgb = bg_color; b.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(l+0.05), Inches(t+0.02), Inches(len(label)*0.11+0.15), Inches(0.3))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = text_color

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

# Decorative circles
for cx, cy, cr, alpha_col in [
    (11.5, 1.0, 2.8, RGBColor(0x00, 0x50, 0x60)),
    (1.5,  6.2, 1.8, RGBColor(0x10, 0x20, 0x55)),
]:
    c = s.shapes.add_shape(9, Inches(cx-cr/2), Inches(cy-cr/2), Inches(cr), Inches(cr))
    c.fill.solid(); c.fill.fore_color.rgb = alpha_col; c.line.fill.background()

# Horizontal teal stripe
bar = s.shapes.add_shape(1, Inches(0), Inches(3.45), Inches(13.33), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

txt(s, "AGASTYA", 1.0, 1.3, 11.0, 1.5, size=72, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, "आगस्त्य", 1.0, 2.55, 11.0, 0.8, size=28, bold=False, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
txt(s, "AI Medication Intelligence System", 1.0, 3.65, 11.0, 0.7, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Knowledge  ·  Care  ·  Intelligence", 1.0, 4.45, 11.0, 0.5, size=16, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)
txt(s, "Built at Hackathon · April 2025", 1.0, 6.5, 11.0, 0.5, size=13, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "The Problem", 0.72, 0.38, 8.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

stats = [
    ("300M+",  "Indians manage multiple daily\nmedications"),
    ("72%",    "Elderly patients miss doses due\nto complex schedules"),
    ("46%",    "Prescription errors from\nhandwriting misreading"),
    ("5+",     "Languages in a single Indian\nhousehold prescription"),
]
for i, (num, label) in enumerate(stats):
    col = i % 2
    row = i // 2
    bx = 0.4 + col * 6.35
    by = 1.4 + row * 2.5
    box(s, bx, by, 5.9, 2.2)
    txt(s, num,   bx+0.3, by+0.2,  5.3, 1.0, size=52, bold=True, color=TEAL)
    txt(s, label, bx+0.3, by+1.15, 5.3, 0.9, size=15, color=LIGHT_GREY)

txt(s, "Agastya bridges this gap with AI — accessible to any age, any language, any literacy level.",
    0.5, 6.7, 12.3, 0.5, size=14, color=GOLD, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is Agastya
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "What is Agastya?", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

txt(s,
    "Agastya is a web-based AI medication management system designed for Indian households — "
    "particularly elderly patients managing multiple medications without a caregiver always present.",
    0.5, 1.25, 12.3, 1.0, size=17, color=WHITE)

features = [
    ("Scan",      "Point camera at any pill bottle or handwritten prescription — Claude AI reads it instantly"),
    ("Schedule",  "Automatic daily Morning / Afternoon / Night schedule with one-tap adherence tracking"),
    ("Dispense",  "Sends commands to a physical IoT pill dispenser over local HTTP"),
    ("Voice",     "Full multilingual voice output — English, Hindi, Tamil, Kannada, Spanish"),
    ("Alerts",    "AI-composed caregiver email alerts via EmailJS when doses are missed"),
    ("Vitals",    "Health score dashboard integrating Samsung Health / Health Connect data"),
]
for i, (title, desc) in enumerate(features):
    col = i % 2; row = i // 2
    bx = 0.4 + col * 6.35; by = 2.45 + row * 1.5
    box(s, bx, by, 6.1, 1.35)
    accent_bar(s, bx, by, h=1.35, color=TEAL if col==0 else GOLD)
    txt(s, title, bx+0.25, by+0.1,  5.6, 0.45, size=16, bold=True, color=TEAL if col==0 else GOLD)
    txt(s, desc,  bx+0.25, by+0.55, 5.6, 0.75, size=12.5, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Tech Stack", 0.72, 0.38, 8.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

layers = [
    ("Frontend",         "React 18 + Vite 5",                                TEAL),
    ("Styling",          "Tailwind CSS v4 + custom CSS variables",            GOLD),
    ("AI / Vision",      "Anthropic Claude API (claude-opus-4-5) — vision + text generation", RGBColor(0xA0,0x6F,0xFF)),
    ("Local Storage",    "IndexedDB via idb — fully offline, no server needed", TEAL),
    ("Email Alerts",     "EmailJS — no backend required",                     GOLD),
    ("Voice Output",     "Web Speech API (speechSynthesis) + espeak-ng fallback", RGBColor(0xFF,0x70,0x70)),
    ("IoT Bridge",       "Python Flask + Flask-CORS proxy on localhost:5000", RGBColor(0x40,0xC0,0x80)),
    ("Hardware",         "IoT pill dispenser (Arduino/RPi) over HTTP",        GOLD),
]
for i, (layer, tech, color) in enumerate(layers):
    col = i % 2; row = i // 2
    bx = 0.4 + col * 6.35; by = 1.3 + row * 1.4
    box(s, bx, by, 6.1, 1.25)
    accent_bar(s, bx+5.97, by, h=1.25, color=color)
    txt(s, layer, bx+0.2, by+0.1,  5.5, 0.45, size=13, bold=True, color=color)
    txt(s, tech,  bx+0.2, by+0.55, 5.5, 0.65, size=12, color=LIGHT_GREY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Medication Scanner (Feature Deep-dive)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Feature: Medication Scanner", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

txt(s, "Scanner.jsx  +  MedAnalysis.jsx  ·  Claude Vision API",
    0.5, 1.2, 12.3, 0.4, size=13, color=GOLD, italic=True)

box(s, 0.4, 1.75, 5.9, 5.4)
txt(s, "HOW IT WORKS", 0.65, 1.9, 5.5, 0.4, size=13, bold=True, color=TEAL)
steps = [
    "1. User points camera at pill bottle or uploads image",
    "2. Image encoded as base64 and sent to Claude Vision API",
    "3. Claude extracts: drug name, dosage, frequency, slot, side effects",
    "4. Cross-checks against patient's current medication list",
    "5. Drug interaction analysis — rated High / Medium / Low severity",
    "6. Generates simplified instructions in patient's preferred language",
    "7. User confirms → medication added to IndexedDB schedule",
]
bullet_box(s, steps, 0.55, 2.4, 5.65, 4.7, size=12.5, gap=0.6)

box(s, 6.7, 1.75, 6.1, 2.5)
txt(s, "Claude API Prompt Strategy", 6.9, 1.9, 5.8, 0.4, size=13, bold=True, color=GOLD)
bullet_box(s, [
    "System role: clinical pharmacist",
    "Returns structured JSON — never free text",
    "Handles blurred / partial labels gracefully",
    "Supports English, Hindi, Tamil, Kannada scripts",
], 6.85, 2.35, 5.75, 2.0, size=12.5, gap=0.46)

box(s, 6.7, 4.5, 6.1, 2.65)
txt(s, "Interaction Detection", 6.9, 4.65, 5.8, 0.4, size=13, bold=True, color=RGBColor(0xFF,0x70,0x70))
bullet_box(s, [
    "Compares new drug against all active medications",
    "Severity scoring: High (red) / Medium (amber) / Low (green)",
    "Warns before adding conflicting medication to schedule",
    "All logic runs inside Claude — no drug-DB needed",
], 6.85, 5.1, 5.75, 2.0, size=12.5, gap=0.46)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Prescription Scanner
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Feature: Prescription Scanner", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

txt(s, "PrescriptionScanner.jsx  ·  PrescriptionLibrary.jsx  ·  PrescriptionDetail.jsx",
    0.5, 1.2, 12.3, 0.4, size=13, color=GOLD, italic=True)

txt(s, "Understands handwritten Indian clinic prescriptions — including mixed regional scripts",
    0.5, 1.7, 12.3, 0.45, size=15, color=WHITE)

cols = [
    ("EXTRACTION", TEAL, [
        "Clinic name & doctor details",
        "Patient name, age, diagnosis",
        "Visit vitals (BP, temp, SpO₂)",
        "All medications with dosage & duration",
        "Prescription date",
    ]),
    ("SHORTHAND UNDERSTOOD", GOLD, [
        "OD = Once daily",
        "BD = Twice daily",
        "TDS = Three times daily",
        "QID = Four times daily",
        "HS = At bedtime  ·  SOS = As needed",
    ]),
    ("SMART DATE CALC", RGBColor(0xA0,0x6F,0xFF), [
        "\"x10 days\" → calculates expiry date",
        "\"5D\" / \"1/52\" / \"2/12\" formats",
        "Auto-removes expired courses",
        "Shows days remaining in schedule",
        "Stored permanently in IndexedDB",
    ]),
]
for i, (title, color, items) in enumerate(cols):
    bx = 0.4 + i * 4.25
    box(s, bx, 2.3, 4.05, 4.9)
    accent_bar(s, bx, 2.3, h=0.06, color=color)
    txt(s, title, bx+0.2, 2.42, 3.7, 0.4, size=12, bold=True, color=color)
    bullet_box(s, items, bx+0.15, 2.92, 3.8, 4.0, size=12.5, dot_color=color, gap=0.52)

txt(s, "Review & correct AI extraction before saving — human stays in the loop.",
    0.5, 7.1, 12.3, 0.4, size=13, color=GOLD, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Multilingual Voice
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Multilingual Voice & Accessibility", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

langs = ["English", "हिन्दी", "தமிழ்", "ಕನ್ನಡ", "Español"]
for i, lang in enumerate(langs):
    bx = 0.4 + i * 2.5
    box(s, bx, 1.3, 2.3, 0.8, fill_color=CARD_BG)
    accent_bar(s, bx, 1.3, w=2.3, h=0.05, color=TEAL)
    txt(s, lang, bx, 1.4, 2.3, 0.6, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

box(s, 0.4, 2.35, 5.9, 4.85)
txt(s, "PHARMACY VOICE MODE", 0.65, 2.5, 5.5, 0.4, size=13, bold=True, color=TEAL)
bullet_box(s, [
    "Full-screen mode — show the pharmacist your phone",
    "Reads out medication name & instructions aloud",
    "Uses patient's preferred language automatically",
    "Web Speech API with smart voice selection",
    "Prefers Google Cloud voices; falls back to espeak-ng",
    "diagnoseVoices() helper for debugging",
], 0.55, 2.95, 5.65, 4.1, size=13, gap=0.58)

box(s, 6.7, 2.35, 6.1, 2.3)
txt(s, "AI TRANSLITERATION", 6.9, 2.5, 5.8, 0.4, size=13, bold=True, color=GOLD)
bullet_box(s, [
    "Claude transliterates medication names into patient's script",
    "Personalised instructions per language + patient age",
    "Works offline (instructions pre-generated at scan time)",
], 6.85, 2.95, 5.75, 1.6, size=12.5, gap=0.5)

box(s, 6.7, 4.9, 6.1, 2.3)
txt(s, "i18n ARCHITECTURE", 6.9, 5.05, 5.8, 0.4, size=13, bold=True, color=RGBColor(0xA0,0x6F,0xFF))
bullet_box(s, [
    "LanguageContext.jsx — React context provider",
    "useT() hook — get translated UI string anywhere",
    "useLang() hook — get/set active language",
    "i18n.js — all 5 languages' string tables",
], 6.85, 5.5, 5.75, 1.65, size=12.5, gap=0.46)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Smart Dispenser + Schedule
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Smart Dispenser & Daily Schedule", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

box(s, 0.4, 1.3, 5.9, 5.9)
txt(s, "DAILY SCHEDULE", 0.65, 1.45, 5.5, 0.4, size=13, bold=True, color=TEAL)
txt(s, "Schedule.jsx", 0.65, 1.9, 5.5, 0.3, size=11, color=GOLD, italic=True)
bullet_box(s, [
    "Three slots: Morning / Afternoon / Night",
    "One-tap mark as taken — live adherence tracking",
    "Shows days remaining for time-limited courses",
    "Auto-expires medications when course ends",
    "Persisted entirely in IndexedDB — no server",
    "Seeded with mockMedications.js on first run",
], 0.55, 2.3, 5.65, 4.6, size=13, gap=0.6)

box(s, 6.7, 1.3, 6.1, 2.85)
txt(s, "IOT DISPENSER BRIDGE", 6.9, 1.45, 5.8, 0.4, size=13, bold=True, color=GOLD)
txt(s, "DispenserBridge.jsx  ·  dispenser.js  ·  Flask server.py", 6.9, 1.88, 5.8, 0.3, size=10.5, color=LIGHT_GREY, italic=True)
bullet_box(s, [
    "HTTP POST to localhost:5000/dispense",
    "Compartment mapping: morning=1, afternoon=2, night=3",
    "Countdown animation before dispensing",
    "Polls /status for confirmation",
    "Graceful offline fallback — app works without dispenser",
], 6.85, 2.3, 5.75, 2.4, size=12.5, gap=0.46)

box(s, 6.7, 4.4, 6.1, 2.8)
txt(s, "ADHERENCE HISTORY", 6.9, 4.55, 5.8, 0.4, size=13, bold=True, color=RGBColor(0xFF,0x70,0x70))
txt(s, "AdherenceHistory.jsx  ·  mockHistory.js", 6.9, 4.98, 5.8, 0.3, size=10.5, color=LIGHT_GREY, italic=True)
bullet_box(s, [
    "30-day log: taken on time / taken late / missed",
    "Summary statistics and visual bar breakdown",
    "Drives the animated health score ring on Dashboard",
    "Seeded with 30 days of realistic mock data",
], 6.85, 5.4, 5.75, 1.7, size=12.5, gap=0.46)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Architecture Overview", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

# Layer boxes
layers_arch = [
    (0.4, 1.3, 12.5, 1.0, TEAL,   "FRONTEND  —  React 18 + Vite 5 + Tailwind CSS v4",
     "14 components  ·  6 utils  ·  LanguageContext  ·  idb IndexedDB"),
    (0.4, 2.55, 5.9, 1.0, GOLD,   "CLAUDE AI  (Anthropic API)",
     "Medication scan  ·  Prescription OCR  ·  Drug interaction  ·  Caregiver alert text  ·  Translation"),
    (6.7, 2.55, 6.2, 1.0, RGBColor(0xA0,0x6F,0xFF), "LOCAL STORAGE  (IndexedDB via idb)",
     "Prescriptions  ·  Medication schedule  ·  Adherence history  ·  Patient profile"),
    (0.4, 3.8,  5.9, 1.0, RGBColor(0x40,0xC0,0x80), "EMAILJS  (Caregiver Alerts)",
     "No backend required  ·  EmailJS SDK called from browser"),
    (6.7, 3.8,  6.2, 1.0, RGBColor(0xFF,0x70,0x70), "WEB SPEECH API  (Voice Output)",
     "speechSynthesis  ·  Smart voice selection  ·  espeak-ng fallback"),
    (0.4, 5.05, 12.5, 1.0, RGBColor(0x30,0x80,0xC0), "OPTIONAL: Flask Bridge  +  IoT Pill Dispenser",
     "Python Flask on localhost:5000  ·  HTTP dispense commands  ·  Graceful offline degradation"),
]
for (l, t, w, h, col, title, sub) in layers_arch:
    box(s, l, t, w, h)
    accent_bar(s, l, t, h=h, color=col)
    txt(s, title, l+0.2, t+0.08, w-0.35, 0.42, size=13, bold=True, color=col)
    txt(s, sub,   l+0.2, t+0.52, w-0.35, 0.45, size=11.5, color=LIGHT_GREY)

txt(s, "Zero-backend by design — runs entirely in the browser. Claude API is the only external call.",
    0.5, 6.3, 12.3, 0.5, size=14, color=GOLD, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Project Structure
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Project Structure", 0.72, 0.38, 8.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

box(s, 0.4, 1.3, 3.85, 5.9)
txt(s, "src/components/  (14 files)", 0.6, 1.42, 3.5, 0.38, size=12, bold=True, color=TEAL)
bullet_box(s, [
    "Dashboard.jsx — health score ring",
    "Scanner.jsx — pill bottle camera",
    "MedAnalysis.jsx — AI results",
    "PrescriptionScanner.jsx",
    "PrescriptionLibrary.jsx",
    "PrescriptionDetail.jsx",
    "Schedule.jsx — daily slots",
    "DispenserBridge.jsx",
    "DispenserSettings.jsx",
    "VitalsPanel.jsx",
    "AdherenceHistory.jsx",
    "PatientProfile.jsx",
    "CaregiverAlert.jsx",
    "PharmacyVoice.jsx",
], 0.5, 1.85, 3.65, 5.2, size=11, gap=0.37)

box(s, 4.5, 1.3, 3.85, 2.75)
txt(s, "src/utils/  (6 files)", 4.7, 1.42, 3.5, 0.38, size=12, bold=True, color=GOLD)
bullet_box(s, [
    "claudeApi.js — all Claude calls",
    "prescriptionDB.js — IndexedDB CRUD",
    "dispenser.js — IoT HTTP client",
    "emailAlert.js — EmailJS sender",
    "voiceEngine.js — Speech API",
    "i18n.js — 5-language strings",
], 4.6, 1.85, 3.65, 2.3, size=11.5, gap=0.38)

box(s, 4.5, 4.3, 3.85, 2.9)
txt(s, "src/data/  (5 mock files)", 4.7, 4.42, 3.5, 0.38, size=12, bold=True, color=RGBColor(0xA0,0x6F,0xFF))
bullet_box(s, [
    "mockPatient.js — default patient",
    "mockMedications.js — sample meds",
    "mockPrescriptions.js — sample Rx",
    "mockVitals.js — reference values",
    "mockHistory.js — 30-day history",
], 4.6, 4.85, 3.65, 2.3, size=11.5, gap=0.38)

box(s, 8.6, 1.3, 4.65, 2.2)
txt(s, "Root config files", 8.8, 1.42, 4.3, 0.38, size=12, bold=True, color=RGBColor(0x40,0xC0,0x80))
bullet_box(s, [
    "vite.config.js — Vite + Tailwind plugin",
    "tailwind.config.js",
    "eslint.config.js",
    "index.html",
    "package.json",
], 8.7, 1.85, 4.45, 1.7, size=11.5, gap=0.38)

box(s, 8.6, 3.75, 4.65, 3.45)
txt(s, "dispenser-bridge/", 8.8, 3.88, 4.3, 0.38, size=12, bold=True, color=RGBColor(0xFF,0x70,0x70))
bullet_box(s, [
    "server.py — Flask proxy + dispenser",
    "Endpoints: /dispense, /status",
    "CORS-enabled for browser access",
    "Runs on localhost:5000 (optional)",
], 8.7, 4.3, 4.45, 1.5, size=11.5, gap=0.42)

txt(s, "Blend files: Dispenser.blend — 3D model of the physical dispenser hardware",
    8.7, 5.5, 4.45, 0.8, size=10.5, color=LIGHT_GREY, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Claude API Integration
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Claude API Integration", 0.72, 0.38, 10.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

txt(s, "All AI calls centralised in  src/utils/claudeApi.js", 0.5, 1.2, 12.3, 0.4, size=13, color=GOLD, italic=True)

calls = [
    ("scanMedication()",       TEAL,                      "Vision + text  ·  Extracts drug, dose, frequency, slot, side-effects from pill bottle image"),
    ("scanPrescription()",     GOLD,                      "Vision + text  ·  Parses handwritten clinic prescription with Indian shorthand"),
    ("checkInteractions()",    RGBColor(0xFF,0x70,0x70),  "Text  ·  Cross-checks new drug against patient's active medication list — returns severity"),
    ("translateInstructions()",RGBColor(0xA0,0x6F,0xFF),  "Text  ·  Translates and transliterates medication instructions into patient's language"),
    ("generateCaregiverAlert()",RGBColor(0x40,0xC0,0x80), "Text  ·  Writes a personalised, compassionate caregiver alert email for missed doses"),
    ("getHealthInsights()",    RGBColor(0xFF,0xC0,0x36),  "Text  ·  Generates contextual health insights from vitals + active medication profile"),
]
for i, (fn, color, desc) in enumerate(calls):
    by = 1.75 + i * 0.9
    box(s, 0.4, by, 12.5, 0.82)
    accent_bar(s, 0.4, by, h=0.82, color=color)
    txt(s, fn,   0.65, by+0.05, 3.8, 0.38, size=14, bold=True, color=color)
    txt(s, desc, 4.6,  by+0.18, 8.2, 0.55, size=12.5, color=LIGHT_GREY)

txt(s, "Model: claude-opus-4-5  ·  All prompts return structured JSON  ·  Images sent as base64 in messages API",
    0.5, 7.1, 12.3, 0.4, size=12.5, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Impact & Future Roadmap
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, 0.5, 0.4, h=0.55)
txt(s, "Impact & Roadmap", 0.72, 0.38, 8.0, 0.65, size=34, bold=True, color=TEAL)
divider(s, 1.12)

box(s, 0.4, 1.3, 5.9, 5.9)
txt(s, "WHO BENEFITS", 0.65, 1.45, 5.5, 0.4, size=13, bold=True, color=TEAL)
bullet_box(s, [
    "Elderly patients managing 5+ medications",
    "Non-English-literate rural households",
    "Single-person households without daily caregiver",
    "Pharmacists serving non-local-language patients",
    "Caregivers monitoring family members remotely",
    "Doctors whose handwritten Rx causes confusion",
], 0.55, 1.9, 5.65, 4.5, size=13.5, gap=0.62)

box(s, 6.7, 1.3, 6.1, 5.9)
txt(s, "WHATS NEXT", 6.9, 1.45, 5.8, 0.4, size=13, bold=True, color=GOLD)
roadmap = [
    "PWA — install to home screen + offline mode",
    "Browser push notifications for reminders",
    "Cloud sync — multi-device via backend",
    "Wearable integration — Fitbit, Apple Watch, Galaxy Watch",
    "Physical dispenser firmware (Arduino / RPi)",
    "Doctor portal — digital prescription signing",
    "Refill reminders + nearby pharmacy locator",
    "Biometric / OTP authentication",
    "Expand languages — Bengali, Telugu, Marathi",
]
bullet_box(s, roadmap, 6.85, 1.9, 5.75, 5.0, size=12.5, dot_color=GOLD, gap=0.56)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)

for cx, cy, cr, col in [
    (6.66, 3.75, 9.0, RGBColor(0x06, 0x14, 0x38)),
    (6.66, 3.75, 6.5, RGBColor(0x08, 0x18, 0x44)),
    (6.66, 3.75, 4.0, RGBColor(0x0C, 0x1E, 0x50)),
]:
    c = s.shapes.add_shape(9, Inches(cx-cr/2), Inches(cy-cr/2), Inches(cr), Inches(cr))
    c.fill.solid(); c.fill.fore_color.rgb = col; c.line.fill.background()

bar = s.shapes.add_shape(1, Inches(0), Inches(3.3), Inches(13.33), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

txt(s, "AGASTYA", 1.0, 1.2, 11.0, 1.4, size=68, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
txt(s, "आगस्त्य", 1.0, 2.55, 11.0, 0.7, size=26, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
txt(s, "Knowledge  ·  Care  ·  Intelligence", 1.0, 3.55, 11.0, 0.55, size=20, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "AI Medication Intelligence — built for every Indian household.", 1.0, 4.3, 11.0, 0.5, size=15, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

for i, (label, val) in enumerate([
    ("Components", "14"), ("Languages", "5"), ("AI Calls", "6"), ("Backend Required", "0"),
]):
    bx = 1.1 + i * 2.8
    box(s, bx, 5.1, 2.5, 1.3)
    txt(s, val,   bx, 5.2, 2.5, 0.7, size=38, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, label, bx, 5.9, 2.5, 0.4, size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

txt(s, "github.com/TGMadhusoodhan/agastya  ·  Hackathon April 2025",
    1.0, 6.75, 11.0, 0.4, size=12, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

# ── Save ───────────────────────────────────────────────────────────────────────
out = "/home/MadhuArch/agastya/Agastya_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}")
